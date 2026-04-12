#!/usr/bin/env python3
"""Binary file extraction: PDF, DOCX, PPTX to markdown.

Extracts text from binary source files and converts to markdown format.
Libraries (pymupdf, python-docx, python-pptx) are optional — the script
raises ImportError at call time if the needed library is missing.

Usage as a library:
    from extract_binary import extract, extract_pdf, extract_docx, extract_pptx
    md = extract("report.pdf")          # auto-detect format
    md = extract_pdf("report.pdf")      # explicit format

Usage as a script:
    python extract_binary.py report.pdf
    python extract_binary.py slides.pptx --format pptx
"""
from __future__ import annotations

import argparse
import sys
from pathlib import Path

# ── Format registry ─────────────────────────────────────────────────

FORMAT_MAP = {
    ".pdf": "pdf",
    ".docx": "docx",
    ".pptx": "pptx",
}

# ── Extraction functions ────────────────────────────────────────────


def extract_pdf(path: str) -> str:
    """Extract text from a PDF file using pymupdf."""
    import fitz  # pymupdf

    doc = fitz.open(path)
    pages = []
    for i, page in enumerate(doc):
        text = page.get_text().strip()
        if text:
            pages.append(text)
    doc.close()
    return "\n\n".join(pages)


def extract_docx(path: str) -> str:
    """Extract text from a DOCX file using python-docx.

    Heading styles are converted to markdown ``#`` headers.
    """
    from docx import Document

    doc = Document(path)
    lines = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        style = para.style.name if para.style else ""
        if style.startswith("Heading"):
            # Extract heading level from style name (e.g. "Heading 1" -> 1)
            try:
                level = int(style.split()[-1])
            except (ValueError, IndexError):
                level = 1
            lines.append(f"{'#' * level} {text}")
        else:
            lines.append(text)
    return "\n\n".join(lines)


def extract_pptx(path: str) -> str:
    """Extract text from a PPTX file using python-pptx.

    Each slide is labeled with ``## Slide N``.
    """
    from pptx import Presentation

    prs = Presentation(path)
    sections = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        texts.append(text)
        slide_text = "\n\n".join(texts) if texts else ""
        sections.append(f"## Slide {i}\n\n{slide_text}")
    return "\n\n".join(sections)


def extract(input_path: str, format: str = None) -> str:
    """Extract text from a binary file to markdown.

    If *format* is None, auto-detects from file extension.
    Raises ``ValueError`` for unsupported formats.
    """
    if format is None:
        ext = Path(input_path).suffix.lower()
        format = FORMAT_MAP.get(ext)
        if format is None:
            raise ValueError(
                f"Unsupported format for '{Path(input_path).name}'. "
                f"Supported extensions: {', '.join(sorted(FORMAT_MAP.keys()))}"
            )

    if format == "pdf":
        return extract_pdf(input_path)
    elif format == "docx":
        return extract_docx(input_path)
    elif format == "pptx":
        return extract_pptx(input_path)
    else:
        raise ValueError(f"Unsupported format: {format}")


# ── CLI ─────────────────────────────────────────────────────────────


def _build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(
        description="Extract text from binary files (PDF, DOCX, PPTX) to markdown."
    )
    parser.add_argument("input_file", help="Path to the binary file")
    parser.add_argument("--format", choices=["pdf", "docx", "pptx"],
                        default=None, help="Force format (auto-detected if omitted)")
    return parser


def main(argv: list[str] = None) -> None:
    parser = _build_parser()
    args = parser.parse_args(argv)
    result = extract(args.input_file, format=args.format)
    print(result)


if __name__ == "__main__":
    main()
